from pptx import Presentation
from io import BytesIO
import zipfile
import re


def is_valid_pptx(file_path):
    try:
        with zipfile.ZipFile(file_path, 'r') as zip_file:
            required_files = ["ppt/presentation.xml", "ppt/slides/slide1.xml"]
            for file in required_files:
                if file not in zip_file.namelist():
                    return False
    except zipfile.BadZipFile:
        return False
    return True

def is_pptx_file(file_path):
    return file_path.lower().endswith('.pptx')

def readFromPowerPoint(first_file, first_file_slide_number, second_file):
    if isinstance(first_file, str):
        if not is_pptx_file(first_file):
            raise ValueError("Invalid file format. Must be a PowerPoint (.pptx) file.")
        if not is_valid_pptx(first_file):
            raise ValueError("Invalid PowerPoint file.")
        with open(first_file, "rb") as file:
            presentation_data = BytesIO(file.read())
    else:
        presentation_data = first_file

    presentation = Presentation(presentation_data)

    if first_file_slide_number <= 0:
        raise ValueError("First file slide input must be greater than 0.")

    first_file_slide = presentation.slides[first_file_slide_number - 1]

    first_file_data = []

    for shape in first_file_slide.shapes:
        if shape.has_table:
            table = shape.table
            table_data = {}
            for row in table.rows:
                attribute_name = row.cells[0].text.strip().replace(':', '')
                attribute_value = row.cells[1].text.strip()
                table_data[attribute_name] = attribute_value
            first_file_data.append({"type": "table", "data": table_data})
        elif hasattr(shape, 'text'):
            first_file_data.append({"type": "text", "data": shape.text})

    
    if isinstance(second_file, str):
        if not is_pptx_file(second_file):
            raise ValueError("Invalid file format. Must be a PowerPoint (.pptx) file.")
        if not is_valid_pptx(second_file):
            raise ValueError("Invalid PowerPoint file.")
        with open(second_file, "rb") as file:
            presentation_data_second = BytesIO(file.read())
    else:
        presentation_data_second = second_file

    presentation_second = Presentation(presentation_data_second)

    second_file_data = []

    for slide in presentation_second.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                table_data = {}
                for row in table.rows:
                    attribute_name = row.cells[0].text.strip().replace(':', '')
                    attribute_value = row.cells[1].text.strip()
                    table_data[attribute_name] = attribute_value
                second_file_data.append({"type": "table", "data": table_data})
            elif hasattr(shape, 'text'):
                second_file_data.append({"type": "text", "data": shape.text})

    return first_file_data, second_file_data

        
def writeInPowerPoint(second_file, first_file_data, second_file_data):
    if not is_pptx_file(second_file) or not is_valid_pptx(second_file):
        raise ValueError("Invalid target file. Must be a valid PowerPoint (.pptx) file.")

    presentation = Presentation(second_file)

    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape_item_second in second_file_data:
            for shape_item_first in first_file_data:
                if shape_item_first["type"] == shape_item_second["type"]:
                    if shape_item_first["type"] == "table":
                        for key, value in shape_item_first["data"].items():
                            if key in shape_item_second['data']:
                                for shape in slide.shapes:
                                    if shape.has_table:
                                        for table_row in shape.table.rows:
                                            for cell in table_row.cells:
                                                if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                                                    old_font_size = cell.text_frame.paragraphs[0].runs[0].font.size
                                                    cell.text = re.sub(r"(?<!\S){}(?!\S)".format(re.escape(shape_item_second['data'][key])), f"{value}", cell.text)
                                                    cell.text_frame.paragraphs[0].runs[0].font.size = old_font_size
                                    elif shape.has_text_frame:
                                        for paragraph in shape.text_frame.paragraphs:
                                            for run in paragraph.runs:
                                                old_font_size = run.font.size
                                                run.text = re.sub(r"(?<!\S){}(?!\S)".format(re.escape(shape_item_second['data'][key])), f"{value}", run.text)
                                                run.font.size = old_font_size
                                                                                                                 
                    
    presentation.save(second_file)






source_pptx_path = r"E:\Projects\mypptx\Presentation1.pptx"
target_pptx_path = r"E:\Projects\mypptx\template1.pptx"
source_slide_number = 1


first_file_data, second_file_data = readFromPowerPoint(source_pptx_path, source_slide_number, target_pptx_path)
writeInPowerPoint(target_pptx_path, first_file_data, second_file_data)
